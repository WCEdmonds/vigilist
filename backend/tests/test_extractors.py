"""Unit tests for the loose-file extraction dispatcher."""

import io

from app.services.extractors import extract, ExtractResult


def _docx_bytes(text: str) -> bytes:
    from docx import Document as Docx
    d = Docx()
    d.add_paragraph(text)
    buf = io.BytesIO()
    d.save(buf)
    return buf.getvalue()


def _xlsx_bytes(value: str) -> bytes:
    from openpyxl import Workbook
    wb = Workbook()
    wb.active["A1"] = value
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _pptx_bytes(text: str) -> bytes:
    from pptx import Presentation
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # title-only layout
    slide.shapes.title.text = text
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def test_extract_docx():
    r = extract("a.docx", _docx_bytes("hello world"))
    assert r.extraction_status == "ok"
    assert "hello world" in r.text
    assert r.file_type == "docx"


def test_extract_xlsx():
    r = extract("b.xlsx", _xlsx_bytes("cell text"))
    assert r.extraction_status == "ok"
    assert "cell text" in r.text


def test_extract_pptx():
    r = extract("c.pptx", _pptx_bytes("slide title"))
    assert r.extraction_status == "ok"
    assert "slide title" in r.text


def test_extract_text_and_case_insensitive_ext():
    r = extract("notes.TXT", b"line one\nline two")
    assert r.extraction_status == "ok"
    assert "line one" in r.text


def test_extract_image_uses_ocr_fn():
    r = extract("scan.png", b"\x89PNG-not-real", ocr_fn=lambda b: "ocr text")
    assert r.extraction_status == "ok"
    assert r.text == "ocr text"
    assert r.file_type == "image"


def test_extract_image_with_real_ocr_jpeg_text_wrapper(monkeypatch):
    """Regression (final review Finding 1): ``_ocr_jpeg`` returns ``PageOcr |
    None`` (the layout-aware contract ``iter_pdf_pages`` needs), but
    ``extract()`` expects a plain ``str`` ocr_fn and calls ``.strip()`` on the
    result — mixing the two contracts raised
    ``'PageOcr' object has no attribute 'strip'`` for every standalone image
    routed through ``extract(ocr_fn=_ocr_jpeg)``.

    This exercises the REAL wiring: monkeypatch the actual Vision entry point
    (``app.services.ocr.ocr_page_vision_bytes``) to return a PageOcr, then
    call ``extract()`` with the real ``_ocr_jpeg_text`` adapter — not a
    hand-stubbed string fn — so a regression that reintroduces the PageOcr
    leak fails with an AttributeError instead of silently passing.
    """
    import app.services.ocr as ocr_mod
    from app.services.ingest_pdf import _ocr_jpeg_text
    from app.services.ocr import PageOcr

    monkeypatch.setattr(ocr_mod, "ocr_page_vision_bytes", lambda b: PageOcr(text="RECOVERED"))

    r = extract("scan.jpg", b"\xff\xd8\xff-fake-jpeg", ocr_fn=_ocr_jpeg_text)

    assert r.text == "RECOVERED"
    assert r.extraction_status == "ok"
    assert r.extraction_status != "error"
    assert r.file_type == "image"


def test_extract_unsupported():
    for name in ("old.doc", "mail.msg", "archive.pst", "weird.xyz", "noext"):
        r = extract(name, b"whatever")
        assert r.extraction_status == "unsupported", name
        assert r.text == ""


def test_extract_corrupt_supported_type_is_error():
    r = extract("broken.docx", b"not a real docx")
    assert r.extraction_status == "error"
    assert r.extraction_error
    assert r.text == ""


def _pptx_direct_text_bytes(text: str) -> bytes:
    """Build a pptx where text is set via text_frame.text (no explicit runs)."""
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2))
    txBox.text_frame.text = text  # sets paragraph text directly, no explicit runs
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def test_extract_pptx_direct_text():
    """para.text must capture text set via text_frame.text (no run objects)."""
    r = extract("direct.pptx", _pptx_direct_text_bytes("Direct Slide Text"))
    assert r.extraction_status == "ok"
    assert "Direct Slide Text" in r.text


def test_extract_corrupt_xlsx_is_error():
    r = extract("broken.xlsx", b"not a real office file")
    assert r.extraction_status == "error"
    assert r.extraction_error
    assert r.text == ""


def test_extract_corrupt_pptx_is_error():
    r = extract("broken.pptx", b"not a real office file")
    assert r.extraction_status == "error"
    assert r.extraction_error
    assert r.text == ""


def _empty_docx_bytes() -> bytes:
    """Build a docx with no paragraphs containing non-whitespace text."""
    from docx import Document as Docx
    d = Docx()
    # Default Docx() has one empty paragraph — leave it as-is.
    buf = io.BytesIO()
    d.save(buf)
    return buf.getvalue()


def test_extract_empty_docx_is_partial():
    r = extract("empty.docx", _empty_docx_bytes())
    assert r.extraction_status == "partial"
    assert r.text == ""
