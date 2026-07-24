"""Unit tests for new extractor formats: POTX, ODT, RTF."""

import io
import zipfile
from unittest.mock import patch

from app.services.extractors import extract, ExtractResult


def _odt_bytes(paragraphs: list[str]) -> bytes:
    """Build a minimal ODT with the given paragraphs in content.xml."""
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w") as zf:
        # ODT requires a mimetype entry (uncompressed)
        zf.writestr("mimetype", "application/vnd.oasis.opendocument.text", compress_type=zipfile.ZIP_STORED)
        # Minimal content.xml with text:p paragraphs
        xml_parts = ['<?xml version="1.0" encoding="UTF-8"?>']
        xml_parts.append('<office:document-content xmlns:text="urn:oasis:names:tc:opendocument:xmlns:text:1.0">')
        for para in paragraphs:
            xml_parts.append(f"<text:p>{para}</text:p>")
        xml_parts.append("</office:document-content>")
        content_xml = "\n".join(xml_parts).encode("utf-8")
        zf.writestr("content.xml", content_xml)
    return buf.getvalue()


def _rtf_bytes(content: str) -> bytes:
    """Encode RTF content as bytes."""
    return content.encode("latin-1", errors="replace")


def test_extract_odt():
    """Test ODT extraction with two paragraphs."""
    odt = _odt_bytes(["First paragraph", "Second paragraph"])
    r = extract("doc.odt", odt)
    assert r.extraction_status == "ok"
    assert r.file_type == "odt"
    assert "First paragraph" in r.text
    assert "Second paragraph" in r.text
    assert "<" not in r.text  # No XML tags in output


def test_extract_odt_empty_is_partial():
    """Test ODT with no meaningful content."""
    odt = _odt_bytes([])
    r = extract("empty.odt", odt)
    assert r.extraction_status == "partial"
    assert r.text == ""


def test_extract_odt_with_nested_tags():
    """Test ODT extraction strips nested tags but preserves paragraph text."""
    odt = _odt_bytes(["Para with <text:span>nested</text:span> tags", "Normal text"])
    r = extract("nested.odt", odt)
    assert r.extraction_status == "ok"
    assert "Para with nested tags" in r.text
    assert "Normal text" in r.text


def test_extract_odt_corrupt_is_error():
    """Test malformed ODT bytes → error status, never raises."""
    r = extract("broken.odt", b"not a real odt")
    assert r.extraction_status == "error"
    assert r.extraction_error
    assert r.text == ""


# --- Finding C1: bounded decompression read for content.xml (bomb-proofing) --
#
# _extract_odt must never buffer an unbounded decompressed content.xml — a
# small, highly-compressible ODT can inflate to a huge in-memory payload.
# _read_zip_member_bounded enforces the cap from actual bytes read (like
# ingest_native's _ZipExploder._read_bounded), independent of any zip header
# claim, and is tested directly here plus via extract()'s error path.


def test_read_zip_member_bounded_stops_at_limit():
    """The bounded zip-member read must never surface more than `limit` bytes
    and must raise (not silently truncate) once the true decompressed size
    exceeds it."""
    import io as _io
    import zipfile as _zipfile
    from app.services.extractors import _read_zip_member_bounded

    payload = b"\x00" * (5 * 1024 * 1024)
    buf = _io.BytesIO()
    with _zipfile.ZipFile(buf, "w", compression=_zipfile.ZIP_DEFLATED) as zf:
        zf.writestr("big.bin", payload)
    zf = _zipfile.ZipFile(_io.BytesIO(buf.getvalue()))

    # A limit far below the true size: raises, never buffers the full payload.
    try:
        _read_zip_member_bounded(zf, "big.bin", limit=1024)
        assert False, "expected ValueError"
    except ValueError:
        pass

    # A limit at/above the true size returns the exact real bytes.
    assert _read_zip_member_bounded(zf, "big.bin", limit=len(payload)) == payload


def test_extract_odt_bomb_guard_errors_without_full_allocation(monkeypatch):
    """A content.xml that is small on disk but decompresses well past the cap
    must never be fully materialized — extract() reports it as an error row,
    mirroring the outer try/except semantics in extract()."""
    import app.services.extractors as extractors_mod

    # Lower the cap so a moderately-sized (but still much-larger-than-cap)
    # payload is enough to trip the guard in a fast test.
    monkeypatch.setattr(extractors_mod, "_ODT_XML_CAP", 1 * 1024 * 1024)

    # Highly compressible content.xml, deflated: small on disk, ~5MB once
    # decompressed — well above the monkeypatched 1MB cap.
    paragraph = "<text:p>" + ("A" * 1000) + "</text:p>"
    xml = '<?xml version="1.0" encoding="UTF-8"?><office:document-content>'
    xml += paragraph * 5000
    xml += "</office:document-content>"
    content_xml = xml.encode("utf-8")

    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w") as zf:
        zf.writestr("mimetype", "application/vnd.oasis.opendocument.text", compress_type=zipfile.ZIP_STORED)
        zf.writestr("content.xml", content_xml, compress_type=zipfile.ZIP_DEFLATED)
    odt = buf.getvalue()
    assert len(odt) < 1 * 1024 * 1024  # the ODT file itself is small (well-compressed)
    assert len(content_xml) > 1 * 1024 * 1024  # but decompresses past the cap

    r = extract("bomb.odt", odt)
    assert r.extraction_status == "error"
    assert r.extraction_error
    assert r.text == ""


def test_extract_potx_routed_to_pptx_extractor():
    """Test .potx reaches _extract_pptx and file_type is 'potx'."""
    # Build a valid PPTX (which is also valid as POTX template format)
    from pptx import Presentation
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Template Slide"
    buf = io.BytesIO()
    prs.save(buf)
    potx_data = buf.getvalue()

    r = extract("template.potx", potx_data)
    assert r.extraction_status == "ok"
    assert r.file_type == "potx"
    assert "Template Slide" in r.text


def test_extract_potx_corrupt_is_error():
    """Test malformed POTX bytes → error status, never raises."""
    r = extract("broken.potx", b"not a real potx")
    assert r.extraction_status == "error"
    assert r.extraction_error
    assert r.text == ""


def test_extract_rtf_strips_control_words():
    """Test RTF extraction removes control words and returns plain text."""
    rtf_content = r"{\rtf1\ansi Hello {\b World}}"
    rtf = _rtf_bytes(rtf_content)
    r = extract("doc.rtf", rtf)
    assert r.extraction_status == "ok"
    assert r.file_type == "rtf"
    assert "Hello" in r.text
    assert "World" in r.text
    assert "\\" not in r.text  # No backslash control words


def test_extract_rtf_empty_is_partial():
    """Test RTF with only control words → partial."""
    rtf_content = r"{\rtf1\ansi}"
    rtf = _rtf_bytes(rtf_content)
    r = extract("empty.rtf", rtf)
    assert r.extraction_status == "partial"


def test_extract_rtf_lenient_never_raises():
    """Test that non-RTF text is handled gracefully (striprtf is lenient)."""
    # striprtf doesn't raise for non-RTF input; it just returns the text as-is
    r = extract("not_rtf.rtf", b"plain text, not rtf")
    assert r.extraction_status in ("ok", "partial")  # text content or empty
    assert r.extraction_error is None  # never crashes


def test_rtf_not_in_text_exts():
    """Test that .rtf is no longer in _TEXT_EXTS by checking control words are stripped."""
    # Build RTF with control sequences that would be visible if passed to _extract_text
    rtf_content = r"{\rtf1\ansi\deff0 {\fonttbl {\f0 Times New Roman;}}Test}"
    rtf = _rtf_bytes(rtf_content)
    r = extract("test.rtf", rtf)
    assert r.extraction_status == "ok"
    # If _extract_text was used, we'd see \ansi, \deff0, \fonttbl etc.
    # striprtf removes them, so verify they're gone
    assert r.text.strip() == "Test"
