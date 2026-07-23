"""Extract text from loose native files using Python libraries.

Pure and deterministic: no DB, no storage. Vision OCR for images is injected
via ``ocr_fn`` so callers/tests control it. Extraction never raises — parse
failures become an ``error`` result so the ingest batch can continue.
"""

from __future__ import annotations

import io
import os
from dataclasses import dataclass


@dataclass
class ExtractResult:
    text: str
    file_type: str
    extraction_status: str            # ok | partial | unsupported | error
    extraction_error: str | None = None


_TEXT_EXTS = {".txt", ".csv", ".md", ".log", ".json", ".xml", ".html", ".htm"}
_IMAGE_EXTS = {".jpg", ".jpeg", ".png", ".tif", ".tiff", ".gif", ".bmp"}
# Extensions we deliberately do not handle here (email → SP4b; legacy binary Office).
# Documented only — the routing catch-all handles these via fall-through.
_UNSUPPORTED_EXTS = {".doc", ".xls", ".ppt", ".msg", ".eml", ".pst"}


def _ext(filename: str) -> str:
    return os.path.splitext(filename or "")[1].lower()


def _status_for(text: str) -> str:
    return "ok" if text.strip() else "partial"


def _extract_docx(data: bytes) -> str:
    from docx import Document as Docx
    doc = Docx(io.BytesIO(data))
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                if cell.text.strip():
                    parts.append(cell.text)
    return "\n".join(parts)


def _extract_xlsx(data: bytes) -> str:
    from openpyxl import load_workbook
    wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    try:
        parts = []
        for ws in wb.worksheets:
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None]
                if cells:
                    parts.append("\t".join(cells))
        return "\n".join(parts)
    finally:
        wb.close()


def _extract_pptx(data: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(data))
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                t = para.text
                if t.strip():
                    parts.append(t)
    return "\n".join(parts)


def _extract_odt(data: bytes) -> str:
    import re
    import zipfile as _zipfile
    with _zipfile.ZipFile(io.BytesIO(data)) as zf:
        xml = zf.read("content.xml").decode("utf-8", errors="replace")
    # <text:p>/<text:h> delimit paragraphs; strip all other tags.
    xml = re.sub(r"</text:(p|h)>", "\n", xml)
    text = re.sub(r"<[^>]+>", "", xml)
    return "\n".join(line.strip() for line in text.splitlines() if line.strip())


def _extract_rtf(data: bytes) -> str:
    from striprtf.striprtf import rtf_to_text
    return rtf_to_text(data.decode("latin-1", errors="replace"), errors="ignore")


def _extract_text(data: bytes) -> str:
    return data.decode("utf-8", errors="replace").replace("\x00", "")


def extract(filename: str, data: bytes, ocr_fn=None) -> ExtractResult:
    """Route ``data`` to the extractor for ``filename``'s extension.

    NOTE: ``.pdf`` is intentionally NOT handled here — callers delegate PDFs to
    ``process_pdf_record`` (page render + OCR). A ``.pdf`` reaching this function
    is treated as unsupported.
    """
    ext = _ext(filename)
    ft = ext.lstrip(".") or "unknown"
    try:
        if ext == ".docx":
            t = _extract_docx(data)
            return ExtractResult(t, "docx", _status_for(t))
        if ext == ".xlsx":
            t = _extract_xlsx(data)
            return ExtractResult(t, "xlsx", _status_for(t))
        if ext in (".pptx", ".potx"):
            t = _extract_pptx(data)
            return ExtractResult(t, ext.lstrip("."), _status_for(t))
        if ext == ".odt":
            t = _extract_odt(data)
            return ExtractResult(t, "odt", _status_for(t))
        if ext == ".rtf":
            t = _extract_rtf(data)
            return ExtractResult(t, "rtf", _status_for(t))
        if ext in _TEXT_EXTS:
            t = _extract_text(data)
            return ExtractResult(t, ft, _status_for(t))
        if ext in _IMAGE_EXTS:
            t = (ocr_fn(data) if ocr_fn else "") or ""
            return ExtractResult(t, "image", _status_for(t))
        # Legacy Office, email, unknown/no extension.
        return ExtractResult("", ft, "unsupported")
    except Exception as e:  # never raise — a bad file is an error row, not a crash
        return ExtractResult("", ft, "error", str(e)[:500])
